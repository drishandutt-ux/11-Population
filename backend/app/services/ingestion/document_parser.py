"""
Document parser — supports:
  Text:          .txt  .md  .rst  .rtf
  Documents:     .pdf  .docx
  Spreadsheets:  .xlsx  .xls  .csv  .tsv  .ods
  Presentations: .pptx
  Data:          .json  .xml  .html  .htm
  Images:        .jpg  .jpeg  .png  .gif  .webp  .bmp  .tiff  .tif  .svg
  Graphics desc: any image format → Claude Vision analysis
"""

from __future__ import annotations
import io
import os
import time

from app.core.monitoring import record_usage_sync


# ---------------------------------------------------------------------------
# Router
# ---------------------------------------------------------------------------

def parse_document(content: bytes, filename: str) -> str:
    lower = filename.lower()

    # ── Plain text variants ─────────────────────────────────────────────────
    if lower.endswith((".txt", ".md", ".rst", ".log")):
        return content.decode("utf-8", errors="replace")

    if lower.endswith(".rtf"):
        return _parse_rtf(content)

    # ── Documents ───────────────────────────────────────────────────────────
    if lower.endswith(".pdf"):
        return _parse_pdf(content)

    if lower.endswith(".docx"):
        return _parse_docx(content)

    # ── Spreadsheets ────────────────────────────────────────────────────────
    if lower.endswith(".xlsx") or lower.endswith(".ods"):
        return _parse_xlsx(content, lower)

    if lower.endswith(".xls"):
        return _parse_xls(content)

    if lower.endswith(".csv"):
        return _parse_csv(content, delimiter=",")

    if lower.endswith(".tsv"):
        return _parse_csv(content, delimiter="\t")

    # ── Presentations ────────────────────────────────────────────────────────
    if lower.endswith(".pptx"):
        return _parse_pptx(content)

    # ── Data / markup ────────────────────────────────────────────────────────
    if lower.endswith(".json"):
        return _parse_json(content)

    if lower.endswith((".xml", ".svg")):
        return _parse_xml(content)

    if lower.endswith((".html", ".htm")):
        return _parse_html(content)

    # ── Images → Claude Vision ───────────────────────────────────────────────
    if lower.endswith((".jpg", ".jpeg", ".png", ".gif", ".webp", ".bmp", ".tiff", ".tif")):
        return _parse_image_vision(content, filename)

    # ── Fallback: try UTF-8 ──────────────────────────────────────────────────
    return content.decode("utf-8", errors="replace")


# ---------------------------------------------------------------------------
# Plain text
# ---------------------------------------------------------------------------

def _parse_rtf(content: bytes) -> str:
    from striprtf.striprtf import rtf_to_text
    return rtf_to_text(content.decode("utf-8", errors="replace"))


# ---------------------------------------------------------------------------
# Documents
# ---------------------------------------------------------------------------

def _parse_pdf(content: bytes) -> str:
    import pypdf
    reader = pypdf.PdfReader(io.BytesIO(content))
    pages = [page.extract_text() or "" for page in reader.pages]
    return "\n\n".join(p for p in pages if p.strip())


def _parse_docx(content: bytes) -> str:
    import docx
    doc = docx.Document(io.BytesIO(content))
    parts: list[str] = []

    # Body paragraphs
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text.strip())

    # Tables
    for table in doc.tables:
        rows = []
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            rows.append(" | ".join(cells))
        if rows:
            parts.append("\n".join(rows))

    return "\n\n".join(parts)


# ---------------------------------------------------------------------------
# Spreadsheets
# ---------------------------------------------------------------------------

def _parse_xlsx(content: bytes, lower: str) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(io.BytesIO(content), read_only=True, data_only=True)
    sheets: list[str] = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows: list[str] = []
        row_count = 0

        for row in ws.iter_rows(values_only=True):
            if all(v is None for v in row):
                continue
            row_count += 1
            if row_count > 2000:
                rows.append(f"… ({ws.max_row - row_count} more rows truncated)")
                break
            rows.append(" | ".join("" if v is None else str(v) for v in row))

        if rows:
            sheets.append(f"Sheet: {sheet_name}\n" + "\n".join(rows))

    wb.close()
    return "\n\n".join(sheets) if sheets else "(empty spreadsheet)"


def _parse_xls(content: bytes) -> str:
    import xlrd
    wb = xlrd.open_workbook(file_contents=content)
    sheets: list[str] = []

    for sheet_name in wb.sheet_names():
        ws = wb.sheet_by_name(sheet_name)
        rows: list[str] = []
        for i in range(min(ws.nrows, 2000)):
            row = ws.row_values(i)
            if any(v not in (None, "") for v in row):
                rows.append(" | ".join("" if v is None else str(v) for v in row))
        if rows:
            sheets.append(f"Sheet: {sheet_name}\n" + "\n".join(rows))

    return "\n\n".join(sheets) if sheets else "(empty spreadsheet)"


def _parse_csv(content: bytes, delimiter: str = ",") -> str:
    import csv
    text = content.decode("utf-8", errors="replace")
    reader = csv.reader(io.StringIO(text), delimiter=delimiter)
    rows: list[str] = []
    for i, row in enumerate(reader):
        if i >= 2000:
            rows.append(f"… ({i} more rows truncated)")
            break
        if any(cell.strip() for cell in row):
            rows.append(" | ".join(row))
    return "\n".join(rows)


# ---------------------------------------------------------------------------
# Presentations
# ---------------------------------------------------------------------------

def _parse_pptx(content: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(content))
    slides: list[str] = []

    for i, slide in enumerate(prs.slides, 1):
        parts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
            # Tables inside slides
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    parts.append(" | ".join(cells))
        if parts:
            slides.append(f"Slide {i}:\n" + "\n".join(parts))

    return "\n\n".join(slides) if slides else "(empty presentation)"


# ---------------------------------------------------------------------------
# Data / markup
# ---------------------------------------------------------------------------

def _parse_json(content: bytes) -> str:
    import json
    try:
        data = json.loads(content.decode("utf-8", errors="replace"))
        return json.dumps(data, indent=2, ensure_ascii=False)
    except Exception:
        return content.decode("utf-8", errors="replace")


def _parse_xml(content: bytes) -> str:
    """Extract all text nodes from XML/SVG."""
    import xml.etree.ElementTree as ET
    try:
        root = ET.fromstring(content)
        texts: list[str] = []
        for elem in root.iter():
            if elem.text and elem.text.strip():
                texts.append(elem.text.strip())
            if elem.tail and elem.tail.strip():
                texts.append(elem.tail.strip())
        return "\n".join(texts) if texts else content.decode("utf-8", errors="replace")
    except Exception:
        return content.decode("utf-8", errors="replace")


def _parse_html(content: bytes) -> str:
    """Strip HTML tags and return plain text."""
    from html.parser import HTMLParser

    class _Extractor(HTMLParser):
        SKIP = {"script", "style", "head", "noscript", "meta", "link"}

        def __init__(self):
            super().__init__()
            self.parts: list[str] = []
            self._skip_depth = 0

        def handle_starttag(self, tag, attrs):
            if tag in self.SKIP:
                self._skip_depth += 1

        def handle_endtag(self, tag):
            if tag in self.SKIP and self._skip_depth:
                self._skip_depth -= 1

        def handle_data(self, data):
            if not self._skip_depth and data.strip():
                self.parts.append(data.strip())

    parser = _Extractor()
    parser.feed(content.decode("utf-8", errors="replace"))
    return " ".join(parser.parts)


# ---------------------------------------------------------------------------
# Images → Claude Vision
# ---------------------------------------------------------------------------

_IMAGE_MEDIA_TYPES = {
    ".jpg": "image/jpeg",
    ".jpeg": "image/jpeg",
    ".png": "image/png",
    ".gif": "image/gif",
    ".webp": "image/webp",
    ".bmp": "image/png",   # Pillow will convert BMP→PNG
    ".tiff": "image/png",  # Pillow will convert TIFF→PNG
    ".tif": "image/png",
}

_VISION_PROMPT = (
    "You are analysing this image as part of a document ingestion pipeline. "
    "Provide a thorough description covering:\n"
    "1. What the image shows (type, subject, context)\n"
    "2. All visible text, labels, headings, captions\n"
    "3. If it is a chart/graph/table: extract every data point, axis label, "
    "legend entry, and trend you can see\n"
    "4. If it is a diagram or infographic: describe each element and its "
    "relationships\n"
    "5. Key figures, numbers, percentages, dates\n"
    "6. Any notable insights or conclusions the image conveys\n"
    "Be exhaustive — this text will be indexed into a knowledge graph."
)


def _parse_image_vision(content: bytes, filename: str) -> str:
    import base64
    import anthropic

    lower = "." + filename.rsplit(".", 1)[-1].lower()
    media_type = _IMAGE_MEDIA_TYPES.get(lower, "image/png")

    # For formats that browsers/Claude don't support natively, convert with Pillow
    if lower in (".bmp", ".tiff", ".tif"):
        content = _convert_to_png(content)

    image_b64 = base64.standard_b64encode(content).decode()

    api_key = os.getenv("ANTHROPIC_API_KEY")
    if not api_key:
        return f"[Image: {filename}] (ANTHROPIC_API_KEY not set — skipped vision analysis)"

    client = anthropic.Anthropic(api_key=api_key)
    _pulse_t0 = time.time()
    message = client.messages.create(
        model="claude-opus-4-5-20251101",
        max_tokens=2048,
        messages=[
            {
                "role": "user",
                "content": [
                    {
                        "type": "image",
                        "source": {
                            "type": "base64",
                            "media_type": media_type,
                            "data": image_b64,
                        },
                    },
                    {"type": "text", "text": _VISION_PROMPT},
                ],
            }
        ],
    )

    record_usage_sync(response=message, model="claude-opus-4-5-20251101", label="doc_vision", started_at=_pulse_t0)
    analysis = message.content[0].text if message.content else "(no response)"
    return f"[Image Analysis — {filename}]\n\n{analysis}"


def _convert_to_png(content: bytes) -> bytes:
    """Use Pillow to convert non-native image formats to PNG bytes."""
    try:
        from PIL import Image
        img = Image.open(io.BytesIO(content)).convert("RGB")
        buf = io.BytesIO()
        img.save(buf, format="PNG")
        return buf.getvalue()
    except Exception:
        return content  # Fall back to original bytes
